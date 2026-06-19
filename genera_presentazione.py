#!/usr/bin/env python3
"""
Genera la presentazione 'Analisi della sosta veicolare a Roma' a partire dai
contenuti di RELAZIONE.md, riusando il formato grafico del template fornito
(Roma Mobilità: banner a strisce rosse, skyline, font Helvetica, accento #C00000).

Uso:  python3 genera_presentazione.py
Out:  Analisi_sosta_veicolare_Roma.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ---------------------------------------------------------------------------
# Percorsi
# ---------------------------------------------------------------------------
ROOT     = os.path.dirname(os.path.abspath(__file__))
# Template grafico Roma Mobilità. Override possibile con la variabile d'ambiente
# PPTX_TEMPLATE; in mancanza si usa il template incluso in presentazione/.
TEMPLATE = os.environ.get(
    "PPTX_TEMPLATE",
    os.path.join(ROOT, "presentazione", "template_roma_mobilita.pptx"))
RESULTS  = os.path.join(ROOT, "results")
MAPS     = os.path.join(RESULTS, "maps")
BRAND    = os.path.join(ROOT, "presentazione", "brand")
BANNER   = os.path.join(BRAND, "banner_red.jpeg")
OUT      = os.path.join(ROOT, "Analisi_sosta_veicolare_Roma.pptx")

# ---------------------------------------------------------------------------
# Palette / stile (dal template)
# ---------------------------------------------------------------------------
RED   = RGBColor(0xC0, 0x00, 0x00)   # accento Roma Mobilità
DARK  = RGBColor(0x30, 0x30, 0x30)   # testo corpo
GRAY  = RGBColor(0x70, 0x70, 0x70)   # testo secondario / caption
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Helvetica"

SW, SH = 14.0, 10.5   # dimensioni slide (pollici)

prs = Presentation(TEMPLATE)
BLANK = prs.slide_layouts[6]   # 'Vuota'


# ---------------------------------------------------------------------------
# Helper
# ---------------------------------------------------------------------------
def _set_run(r, size, bold=False, color=DARK, italic=False, font=FONT):
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def add_banner(slide):
    """Banner a strisce rosse in alto a sinistra (cifra grafica del template)."""
    slide.shapes.add_picture(BANNER, Inches(0.31), Inches(0.0),
                             width=Inches(3.0), height=Inches(0.45))


def add_title(slide, text, subtitle=None):
    """Titolo di sezione: rosso bold Helvetica."""
    tb = slide.shapes.add_textbox(Inches(0.53), Inches(0.62),
                                  Inches(SW - 1.0), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    _set_run(r, 28, bold=True, color=RED)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        _set_run(r2, 15, bold=False, color=GRAY, italic=True)
    # linea rossa sottile sotto al titolo
    line = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0.55), Inches(1.52), Inches(SW - 1.1), Pt(2.2))
    line.fill.solid(); line.fill.fore_color.rgb = RED
    line.line.fill.background()
    line.shadow.inherit = False
    return tb


def add_bullets(slide, items, left, top, width, height,
                size=18, space_after=10):
    """
    items: lista di bullet. Ogni bullet è:
      - str  -> testo normale
      - list di (testo, stile) con stile in {'n','b','h'}:
            n = normale (#303030), b = bold rosso, h = bold scuro
      - dict {'segs': [...], 'level': int, 'bullet': bool, 'size': int}
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        level = 0
        bullet = True
        bsize = size
        if isinstance(it, dict):
            segs = it.get("segs", [])
            level = it.get("level", 0)
            bullet = it.get("bullet", True)
            bsize = it.get("size", size)
        elif isinstance(it, str):
            segs = [(it, 'n')]
        else:
            segs = it

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(space_after)
        p.line_spacing = 1.04

        prefix = ("•   " if level == 0 else "–   ") if bullet else ""
        if prefix:
            rp = p.add_run(); rp.text = prefix
            _set_run(rp, bsize, bold=False,
                     color=RED if level == 0 else GRAY)
        for txt, st in segs:
            r = p.add_run(); r.text = txt
            if st == 'b':
                _set_run(r, bsize, bold=True, color=RED)
            elif st == 'h':
                _set_run(r, bsize, bold=True, color=DARK)
            else:
                _set_run(r, bsize, bold=False, color=DARK)
    return tb


def _img_ratio(path):
    with Image.open(path) as im:
        w, h = im.size
    return w / h


def add_image_fit(slide, path, box_left, box_top, box_w, box_h,
                  align="center", valign="middle"):
    """Inserisce l'immagine adattandola (contain) al box dato, preservando l'aspetto."""
    ar = _img_ratio(path)
    box_ar = box_w / box_h
    if ar >= box_ar:          # immagine più larga -> vincolo sulla larghezza
        w = box_w; h = box_w / ar
    else:                     # più alta -> vincolo sull'altezza
        h = box_h; w = box_h * ar
    if align == "center":
        left = box_left + (box_w - w) / 2
    elif align == "left":
        left = box_left
    else:
        left = box_left + (box_w - w)
    if valign == "middle":
        top = box_top + (box_h - h) / 2
    elif valign == "top":
        top = box_top
    else:
        top = box_top + (box_h - h)
    return slide.shapes.add_picture(path, Inches(left), Inches(top),
                                    width=Inches(w), height=Inches(h))


def add_caption(slide, text, left, top, width, size=12, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _set_run(r, size, italic=True, color=GRAY)
    return tb


def add_kpi_row(slide, kpis, top=2.0, left=0.7, total_w=SW - 1.4, height=1.7):
    """Riga di 'KPI card': lista di (valore, etichetta)."""
    n = len(kpis)
    gap = 0.3
    card_w = (total_w - gap * (n - 1)) / n
    for i, (val, lab) in enumerate(kpis):
        x = left + i * (card_w + gap)
        card = slide.shapes.add_shape(5,  # rounded rectangle
                                      Inches(x), Inches(top),
                                      Inches(card_w), Inches(height))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xF0)
        card.line.color.rgb = RED; card.line.width = Pt(1.25)
        card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Pt(6); tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        _set_run(r, 30, bold=True, color=RED)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = lab
        _set_run(r2, 12.5, color=DARK)


def add_table(slide, data, left, top, width, height,
              header=True, col_widths=None, fsize=13):
    """data: lista di righe (liste di stringhe). Prima riga = header se header=True."""
    rows = len(data); cols = len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                Inches(width), Inches(height)).table
    if col_widths:
        for j, cw in enumerate(col_widths):
            gt.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = gt.cell(i, j)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            if header and i == 0:
                _set_run(r, fsize, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = RED
            else:
                _set_run(r, fsize, color=DARK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (WHITE if i % 2 else
                                            RGBColor(0xF5, 0xF0, 0xF0))
    return gt


def new_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    add_banner(s)
    if title:
        add_title(s, title, subtitle)
    return s


def add_pagenum(slide, n):
    tb = slide.shapes.add_textbox(Inches(SW - 1.2), Inches(SH - 0.5),
                                  Inches(0.9), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(n)
    _set_run(r, 11, color=GRAY)


# ---------------------------------------------------------------------------
# 1. TITOLO  — riuso slide 0 del template, aggiungo sottotitolo
# ---------------------------------------------------------------------------
title_slide = prs.slides[0]
sub = title_slide.shapes.add_textbox(Inches(1.65), Inches(4.75),
                                     Inches(11.09), Inches(1.4))
tf = sub.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Dove, quando e quanto sostano i veicoli nel Comune di Roma"
_set_run(r, 22, bold=False, color=DARK)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Analisi su dati Floating Car Data (FCD) — marzo 2023"
_set_run(r2, 16, italic=True, color=GRAY)

# Rimuovo la slide 1 del template (contenuto 'black point', non pertinente).
# Va eliminata anche la relazione, altrimenti la part resta nel package e
# genera collisioni di nome con le slide aggiunte dopo.
def _delete_slide(prs, index):
    rels_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    sldId = list(prs.slides._sldIdLst)[index]
    rId = sldId.get(rels_ns)
    prs.part.drop_rel(rId)                 # rimuove relazione (e orfana la part)
    prs.slides._sldIdLst.remove(sldId)     # rimuove dall'elenco slide

_delete_slide(prs, 1)


# ===========================================================================
# 2. OBIETTIVO E DATI
# ===========================================================================
s = new_slide("Obiettivo dell'analisi e dati di partenza")
add_bullets(s, [
    [("Tre domande sulla ", 'n'), ("sosta veicolare a Roma", 'b'),
     (", a partire da dati GPS reali:", 'n')],
    {"segs": [("Dove", 'b'), (" parcheggiano i veicoli — su strada (on-street) o fuori strada (off-street)", 'n')], "level": 1},
    {"segs": [("Quanti", 'b'), (" veicoli sono fermi in un dato momento rispetto a quelli in uso", 'n')], "level": 1},
    {"segs": [("Quanto spazio", 'b'), (" stradale è occupato dalla sosta", 'n')], "level": 1},
], 0.7, 1.85, 7.5, 4.2, size=18)

# dati di partenza (riquadro a destra)
box = s.shapes.add_shape(5, Inches(8.5), Inches(1.95), Inches(4.85), Inches(4.6))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xF0)
box.line.color.rgb = RED; box.line.width = Pt(1.25); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_right = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Dati di partenza"
_set_run(r, 17, bold=True, color=RED)
for head, body in [
    ("FCD — spostamenti", "411.356 viaggi GPS (marzo 2023); la destinazione = luogo di sosta"),
    ("municipi_2013.shp", "15 poligoni → perimetro del Comune di Roma"),
    ("AC_VEI.shp", "79.085 poligoni di superficie carrabile (~135,1 km²)"),
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    r = p.add_run(); r.text = head
    _set_run(r, 14, bold=True, color=DARK)
    p = tf.add_paragraph()
    r = p.add_run(); r.text = body
    _set_run(r, 13, color=DARK)
add_pagenum(s, 2)


# ===========================================================================
# 3. METODOLOGIA (pipeline)
# ===========================================================================
s = new_slide("Metodologia: la pipeline di analisi",
              "Sistema di riferimento metrico EPSG:25833 — script Python (geopandas, pandas, matplotlib, h3)")
steps = [
    ("1 · 2", "Caricamento & georeferenziazione", "Unione dei 4 CSV, parsing orari, riproiezione delle destinazioni in coordinate metriche"),
    ("3", "Filtro sul Comune di Roma", "Spatial join (R-tree) dei punti sul perimetro municipale"),
    ("4", "On-street / off-street", "Test dei punti sulla superficie carrabile con buffer 0–10 m"),
    ("5 · 6", "Analisi temporale", "Aggregazione per ora del giorno e giorno della settimana"),
    ("7", "Fermi vs in uso", "Snapshot orari: veicoli in viaggio vs parcheggiati"),
    ("9 · 10", "Scala città & mappe", "Calibrazione sul traffico reale e mappe H3 per municipio"),
]
n = len(steps); gap = 0.25
cw = (SW - 1.4 - gap * (n - 1)) / n
for i, (num, ttl, desc) in enumerate(steps):
    x = 0.7 + i * (cw + gap)
    card = s.shapes.add_shape(5, Inches(x), Inches(2.3), Inches(cw), Inches(4.6))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RED; card.line.width = Pt(1.5); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(10)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "STEP"; _set_run(r, 10, bold=True, color=GRAY)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; _set_run(r, 26, bold=True, color=RED)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
    r = p.add_run(); r.text = ttl; _set_run(r, 13.5, bold=True, color=DARK)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(6)
    r = p.add_run(); r.text = desc; _set_run(r, 11.5, color=DARK)
add_pagenum(s, 3)


# ===========================================================================
# 4. FILTRO TERRITORIALE  (KPI)
# ===========================================================================
s = new_slide("Filtraggio territoriale: gli spostamenti dentro Roma")
add_kpi_row(s, [
    ("411.356", "Spostamenti FCD totali"),
    ("357.117", "Destinazione dentro Roma"),
    ("86,8%", "Quota interna al Comune"),
], top=2.1)
add_bullets(s, [
    [("Ogni destinazione è confrontata con il perimetro municipale tramite ", 'n'),
     ("spatial join indicizzato (R-tree)", 'b'), (".", 'n')],
    [("L'", 'n'), ("86,8%", 'b'), (" degli spostamenti (357.117) ha destinazione interna a Roma; "
      "il restante 13,2% (54.239) è escluso dall'analisi.", 'n')],
    [("Tutte le analisi successive lavorano sui ", 'n'),
     ("357.117 spostamenti", 'h'), (" con destinazione comunale.", 'n')],
], 0.9, 4.4, 12.2, 3.0, size=18, space_after=14)
add_pagenum(s, 4)


# ===========================================================================
# 5. DOVE SI PARCHEGGIA  — on/off-street + buffer  (testo + grafico)
# ===========================================================================
s = new_slide("Dove si parcheggia: sosta su strada vs fuori strada")
add_bullets(s, [
    [("Senza tolleranza il ", 'n'), ("65%", 'b'),
     (" delle soste è esattamente sulla carreggiata (on-street).", 'n')],
    [("Con appena ", 'n'), ("2 m di buffer", 'b'),
     (" — coerenti con l'errore GPS urbano — si sale al ", 'n'), ("76,4%", 'b'),
     (", e a 5 m all'84%.", 'n')],
    [("La curva cresce ripida nei primi metri: i punti 'fuori strada' sono in gran "
      "parte soste su strada ", 'n'), ("mal georeferenziate", 'h'),
     (", non aree private.", 'n')],
    [("Risultato robusto: a Roma si parcheggia ", 'n'),
     ("prevalentemente su strada", 'b'), (".", 'n')],
], 0.7, 1.85, 6.4, 4.6, size=17, space_after=12)
# tabella buffer
add_table(s, [
    ["Buffer", "On-street", "% On"],
    ["0 m", "232.071", "65,0%"],
    ["2 m", "272.941", "76,4%"],
    ["5 m", "300.996", "84,3%"],
    ["10 m", "321.846", "90,1%"],
], left=0.9, top=6.7, width=5.8, height=2.6, col_widths=[1.9, 2.3, 1.6], fsize=13)
add_image_fit(s, os.path.join(RESULTS, "fig4_curva_buffer.png"),
              7.4, 1.9, 6.2, 5.0)
add_caption(s, "Fig. 4 — Quota on-street al variare del buffer sulla superficie carrabile",
            7.4, 7.0, 6.2)
add_pagenum(s, 5)


# ===========================================================================
# 6. ANDAMENTO TEMPORALE DEGLI ARRIVI  (immagine grande)
# ===========================================================================
s = new_slide("Quando si arriva: profilo orario e settimanale")
add_image_fit(s, os.path.join(RESULTS, "fig3_heatmap_volume_ora_giorno.png"),
              0.6, 1.9, 7.6, 5.6, align="center")
add_caption(s, "Fig. 3 — Heatmap del volume di arrivi per ora e giorno",
            0.6, 7.55, 7.6)
add_bullets(s, [
    [("Profilo ", 'n'), ("bimodale", 'b'),
     (": punta mattutina (8–12, ~24.000/h) e serale (~18, ~25.000/h).", 'n')],
    [("Salita mattutina rapida: tra le 6 e le 8 il volume ", 'n'),
     ("quadruplica", 'h'), (".", 'n')],
    [("Netto calo nel ", 'n'), ("weekend", 'b'),
     (": la domenica (~30.000) è meno della metà del picco feriale.", 'n')],
    [("Il 'blocco caldo' feriale 8–19 è ben delimitato; sabato e domenica più "
      "'freddi' e spostati sulle ore centrali.", 'n')],
], 8.5, 2.2, 4.9, 5.0, size=16, space_after=14)
add_pagenum(s, 6)


# ===========================================================================
# 7. ON/OFF-STREET NEL TEMPO
# ===========================================================================
s = new_slide("La sosta su strada cambia tra giorno e notte")
add_image_fit(s, os.path.join(RESULTS, "fig6_andamento_orario_onoff.png"),
              0.6, 1.9, 7.8, 5.7)
add_caption(s, "Fig. 6 — Volume arrivi (sopra) e % on-street per ora (sotto)",
            0.6, 7.65, 7.8)
add_bullets(s, [
    [("La quota on-street ", 'n'), ("non è costante", 'b'), (" nelle 24 ore.", 'n')],
    [("Massima di notte e prima mattina (", 'n'), ("~72–73% tra le 3 e le 5", 'b'),
     ("): i residenti senza box restano in strada.", 'n')],
    [("Minima nelle ore centrali/serali (", 'n'), ("~62–64%", 'h'),
     ("): pesano le destinazioni con parcheggi privati.", 'n')],
    [("La quota on-street è invece ", 'n'), ("stabile lungo la settimana", 'b'),
     (": tratto strutturale del territorio.", 'n')],
], 8.7, 2.3, 4.8, 5.0, size=16, space_after=16)
add_pagenum(s, 7)


# ===========================================================================
# 8. FERMI VS IN USO  (KPI + grafico)
# ===========================================================================
s = new_slide("Un'auto è un oggetto prevalentemente fermo")
add_kpi_row(s, [
    ("94,2%", "Parcheggiati — media 24h"),
    ("98,6%", "Parcheggiati — notte (23–05)"),
    ("~90%", "Fermi anche nell'ora di punta"),
], top=2.0, height=1.6)
add_image_fit(s, os.path.join(RESULTS, "fig9_pct_parcheggiati_vs_inuso_orario.png"),
              0.7, 3.9, 7.6, 3.9, align="left")
add_caption(s, "Fig. 9 — % veicoli parcheggiati vs in uso per ora del giorno",
            0.7, 7.85, 7.6)
add_bullets(s, [
    [("Metodo a ", 'n'), ("snapshot orari", 'b'),
     (" (istantanea alle :30) su 23.571 veicoli, 32 giorni.", 'n')],
    [("Anche nell'ora di massimo traffico, solo ", 'n'), ("~10%", 'b'),
     (" dei veicoli è in movimento.", 'n')],
    [("La variabilità è quasi solo ", 'n'), ("oraria", 'h'),
     (" (giorno/notte), non settimanale.", 'n')],
    [("Lo spazio richiesto dalla sosta è ", 'n'), ("permanente", 'b'),
     (", non di picco.", 'n')],
], 8.6, 4.0, 4.9, 3.8, size=16, space_after=14)
add_pagenum(s, 8)


# ===========================================================================
# 9. STIMA A SCALA DI CITTÀ  (calibrazione + tabella + grafico)
# ===========================================================================
s = new_slide("Dalla scala campione alla scala città",
              "Calibrazione su un riferimento reale di traffico di punta")
add_bullets(s, [
    [("Riferimento esterno: ", 'n'), ("300.000 veicoli unici", 'b'),
     (" in circolazione nella punta 07–09.", 'n')],
    [("Nel campione: ", 'n'), ("1.381 veicoli/giorno", 'n'),
     (" nella stessa finestra → fattore di espansione ", 'n'), ("≈ 217×", 'b'), (".", 'n')],
    [("Superficie occupata = parco totale (", 'n'), ("1.600.000", 'h'),
     (") − veicoli in movimento, con ", 'n'), ("12,5 m²/veicolo", 'h'), (".", 'n')],
], 0.7, 1.95, 6.5, 2.7, size=16, space_after=11)
add_table(s, [
    ["Fascia", "Parcheggiati", "% parch.", "On-street", "% sup."],
    ["Media 24h", "1.538.612", "96,2%", "1.175.500", "10,9%"],
    ["Diurno 06–22", "1.518.292", "94,9%", "1.159.975", "10,7%"],
    ["Notturno 23–05", "1.587.960", "99,2%", "1.213.202", "11,2%"],
    ["Picco 8:30", "1.487.111", "92,9%", "1.136.153", "10,5%"],
], left=0.7, top=4.9, width=6.6, height=2.6,
   col_widths=[1.7, 1.5, 1.0, 1.4, 1.0], fsize=12)
add_bullets(s, [
    [("≈ ", 'n'), ("1,5 milioni di veicoli fermi", 'b'),
     (" in un momento medio, di cui ", 'n'), ("~1,18 mln su strada", 'b'),
     (" ≈ ", 'n'), ("11% della superficie carrabile", 'b'), (".", 'n')],
], 0.7, 7.7, 6.6, 1.6, size=14, space_after=6)
add_image_fit(s, os.path.join(RESULTS, "fig14_stima_citta_fasce_orarie.png"),
              7.6, 1.95, 6.0, 5.6)
add_caption(s, "Fig. 14 — Distribuzione dei veicoli per fascia oraria (stima città)",
            7.6, 7.65, 6.0)
add_pagenum(s, 9)


# ===========================================================================
# 10. MAPPE — dove si concentra la sosta (densità + residenziale)
# ===========================================================================
s = new_slide("Geografia della sosta: densità e funzione residenziale",
              "Griglia esagonale H3 (resolution 9, lato ~174 m)")
add_image_fit(s, os.path.join(MAPS, "map01_densita_totale_campione.png"),
              0.5, 1.95, 6.4, 5.3)
add_caption(s, "Map 01 — Densità di sosta totale (soste/km², campione)",
            0.5, 7.35, 6.4)
add_image_fit(s, os.path.join(MAPS, "map07_indice_residenziale.png"),
              7.1, 1.95, 6.4, 5.3)
add_caption(s, "Map 07 — Indice residenziale (% soste notturne)",
            7.1, 7.35, 6.4)
add_bullets(s, [
    [("La sosta si addensa lungo gli assi radiali e nel ", 'n'),
     ("quadrante sud-est", 'b'),
     (" (Appio–Tuscolano, Cinecittà). Il centro mostra la maggiore ", 'n'),
     ("impronta residenziale", 'b'), (" (% notturna più alta).", 'n')],
], 0.6, 7.8, 12.9, 1.2, size=15, space_after=4)
add_pagenum(s, 10)


# ===========================================================================
# 11. MAPPE PER MUNICIPIO  (barchart comparativo)
# ===========================================================================
s = new_slide("Confronto tra municipi: densità e residenzialità")
add_image_fit(s, os.path.join(MAPS, "map12_municipio_barchart.png"),
              0.6, 1.9, 7.9, 5.7)
add_caption(s, "Map 12 — Indicatori di sosta per municipio",
            0.6, 7.65, 7.9)
add_bullets(s, [
    [("Netta gerarchia di ", 'n'), ("densità", 'b'),
     (": Municipi ", 'n'), ("IX (~1.000 soste/km²)", 'h'),
     (" e ", 'n'), ("VIII (~999)", 'h'), (" staccano tutti.", 'n')],
    [("La ", 'n'), ("residenzialità", 'b'),
     (" segue un ordine quasi opposto: massima nel ", 'n'),
     ("Centro (Municipio I, 23,9%)", 'h'), (".", 'n')],
    [("Disaccoppiamento netto: ", 'n'),
     ("'tanta sosta' ≠ 'sosta da residenti'", 'b'), (".", 'n')],
    [("Due politiche diverse: ", 'n'),
     ("rotazione/tariffazione", 'h'), (" vs ", 'n'),
     ("permessi residenti", 'h'), (".", 'n')],
], 8.8, 2.2, 4.7, 5.2, size=16, space_after=16)
add_pagenum(s, 11)


# ===========================================================================
# 12. SATURAZIONE  (map14 + nota critica)
# ===========================================================================
s = new_slide("Saturazione: i cluster di massima pressione",
              "Indicatore RELATIVO — non leggere i valori assoluti come % di superficie")
add_image_fit(s, os.path.join(MAPS, "map14_saturazione_h3.png"),
              0.5, 1.95, 7.3, 5.5)
add_caption(s, "Map 14 — Saturazione per cella H3",
            0.5, 7.55, 7.3)
# riquadro nota critica
box = s.shapes.add_shape(5, Inches(8.1), Inches(2.1), Inches(5.4), Inches(5.2))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFB, 0xEC, 0xEC)
box.line.color.rgb = RED; box.line.width = Pt(1.5); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_right = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "⚠  Nota di lettura"
_set_run(r, 16, bold=True, color=RED)
for seg in [
    [("I valori superano il 100% (fino a ~2.700% nel Mun. VIII) e ", 'n'),
     ("non sono % di superficie occupata", 'h'), (".", 'n')],
    [("Sommano ", 'n'), ("eventi di sosta cumulati sul mese", 'h'),
     (" (rotazione degli stalli), non occupazioni simultanee.", 'n')],
    [("Vanno usati solo per confrontare la ", 'n'),
     ("pressione relativa tra zone", 'b'), (".", 'n')],
    [("I cluster rossi sono compatti nel ", 'n'),
     ("sud-est (Mun. VII–VIII–IX)", 'b'), (".", 'n')],
    [("La stima fisica affidabile resta l'", 'n'),
     ("~11% di superficie", 'b'), (" (snapshot istantaneo, STEP 9).", 'n')],
]:
    p = tf.add_paragraph(); p.space_before = Pt(9); p.line_spacing = 1.03
    rp = p.add_run(); rp.text = "•  "; _set_run(rp, 14, color=RED)
    for txt, st in seg:
        r = p.add_run(); r.text = txt
        if st == 'b': _set_run(r, 14, bold=True, color=RED)
        elif st == 'h': _set_run(r, 14, bold=True, color=DARK)
        else: _set_run(r, 14, color=DARK)
add_pagenum(s, 12)


# ===========================================================================
# 13. LIMITI DEL METODO
# ===========================================================================
s = new_slide("Limiti del metodo: stime di ordine di grandezza")
add_bullets(s, [
    [("Campione FCD", 'b'), (" non casuale (veicoli recenti, uso intenso) e su un ", 'n'),
     ("solo mese", 'h'), (" (marzo 2023).", 'n')],
    [("Calibrazione", 'b'), (": il fattore 217× confronta grandezze non perfettamente omogenee.", 'n')],
    [("Parco ACI (1,6 mln)", 'b'), (" include mezzi fermi/non assicurati → sovrastima i parcheggiati.", 'n')],
    [("Ratio on/off-street", 'b'), (" estesa a tutta la flotta e costante 24h → sovrastima l'on-street.", 'n')],
    [("Errore GPS", 'b'), (" 5–15 m; nessuna soglia di buffer universalmente corretta.", 'n')],
    [("Ingombro 12,5 m²", 'b'), (" tarato sull'auto media: moto e furgoni hanno ingombri diversi.", 'n')],
    [("AC_VEI ~2014", 'b'), (": trasformazioni viabilistiche successive non incluse.", 'n')],
], 0.8, 1.9, 8.2, 5.4, size=16, space_after=11)
# riquadro direzione distorsione
box = s.shapes.add_shape(5, Inches(9.3), Inches(2.1), Inches(4.2), Inches(4.9))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xF0)
box.line.color.rgb = RED; box.line.width = Pt(1.25); box.shadow.inherit = False
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_right = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Direzione complessiva"
_set_run(r, 16, bold=True, color=RED)
p = tf.add_paragraph(); p.space_before = Pt(10)
r = p.add_run(); r.text = "Le distorsioni puntano a una "
_set_run(r, 14, color=DARK)
r = p.add_run(); r.text = "leggera sovrastima"
_set_run(r, 14, bold=True, color=RED)
r = p.add_run(); r.text = " della sosta su strada."
_set_run(r, 14, color=DARK)
p = tf.add_paragraph(); p.space_before = Pt(12)
r = p.add_run(); r.text = "I valori assoluti = limite superiore ragionevole."
_set_run(r, 14, color=DARK)
p = tf.add_paragraph(); p.space_before = Pt(12); p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Incertezza ±15–25%"
_set_run(r, 19, bold=True, color=RED)
add_pagenum(s, 13)


# ===========================================================================
# 14. CONCLUSIONI
# ===========================================================================
s = new_slide("Sintesi conclusiva")
add_kpi_row(s, [
    ("65% → 76%", "Sosta su strada (0 m → 2 m)"),
    (">94%", "Veicoli fermi nel momento medio"),
    ("~1,5 mln", "Veicoli parcheggiati a Roma"),
    ("~11%", "Superficie carrabile occupata"),
], top=2.0, height=1.65)
add_bullets(s, [
    [("L'", 'n'), ("86,8%", 'b'), (" degli spostamenti ha destinazione dentro Roma; la sosta è ", 'n'),
     ("prevalentemente su strada", 'b'), (" (65% a 0 m, 76% a 2 m).", 'n')],
    [("In un momento medio ", 'n'), (">94%", 'b'),
     (" dei veicoli è fermo (notte >98%): ", 'n'),
     ("lo spazio di sosta è una domanda permanente", 'h'), (".", 'n')],
    [("Geografia: pressione massima nel ", 'n'),
     ("quadrante sud-est (Mun. VII–VIII–IX)", 'b'),
     ("; massima impronta residenziale nel ", 'n'),
     ("Centro Storico (24% notturne)", 'h'), ("; on-street uniforme (73–81%).", 'n')],
    [("Le mappe di ", 'n'), ("saturazione", 'b'),
     (" sono solo indicatori relativi; la stima affidabile di occupazione è l'", 'n'),
     ("~11%", 'b'), (".", 'n')],
    [("Tutte le stime assolute sono ", 'n'), ("ordini di grandezza", 'h'),
     (" con incertezza ±15–25%.", 'n')],
], 0.8, 4.2, 12.6, 4.6, size=17, space_after=16)
add_pagenum(s, 14)


# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"Presentazione salvata: {OUT}")
print(f"Slide totali: {len(prs.slides.__iter__.__self__._sldIdLst)}")
